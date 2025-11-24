"""
PowerPoint operations for CV to PPT conversion
Handles reading sample PPT, understanding structure, and filling data
"""
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime


def _format_left_box_content(area_of_expertise, education):
    """Format left box content with Areas of Expertise and Education sections
    
    Args:
        area_of_expertise: Comma-separated technical skills string
        education: Education details (may contain newlines for multiple entries)
    
    Returns:
        Formatted text string with headers and bullet points
    """
    formatted_parts = []
    
    # Format Areas of Expertise section
    if area_of_expertise and area_of_expertise.strip() and area_of_expertise != 'Not Found':
        # Split by comma and clean up
        skills = [skill.strip() for skill in area_of_expertise.split(',') if skill.strip()]
        
        if skills:
            formatted_parts.append("Areas of Expertise")
            # Add each skill as a bullet point
            for skill in skills:
                formatted_parts.append(f"• {skill}")
    
    # Format Education section (only if education exists)
    if education and education.strip() and education != 'Not Found':
        # Split by newline and clean up
        education_entries = [edu.strip() for edu in education.split('\n') if edu.strip()]
        
        # Also check for | separator (fallback)
        if not education_entries:
            education_entries = [edu.strip() for edu in education.split('|') if edu.strip()]
        
        if education_entries:
            # Add spacing before Education section if Areas of Expertise exists
            if formatted_parts:
                formatted_parts.append("")  # Empty line separator
            
            formatted_parts.append("Education")
            # Add only the highest/most recent education (first entry)
            formatted_parts.append(f"• {education_entries[0]}")
    
    # Join all parts with newlines
    return '\n'.join(formatted_parts)


def _replace_text_with_bold_headers(text_frame, formatted_text):
    """Replace text in text_frame with formatted text, making headers bold
    
    Args:
        text_frame: TextFrame object from python-pptx
        formatted_text: Text string with headers and bullet points
    """
    if not text_frame.paragraphs:
        text_frame.add_paragraph()
    
    # Get formatting from first paragraph (preserve font, size, style)
    first_para = text_frame.paragraphs[0]
    original_runs = first_para.runs
    
    font_name = None
    font_size = None
    font_bold = None
    font_italic = None
    font_color_rgb = None
    para_alignment = None
    para_level = 0
    
    if original_runs:
        first_run = original_runs[0]
        try:
            font_name = first_run.font.name if first_run.font.name else None
        except:
            pass
        try:
            font_size = first_run.font.size
        except:
            pass
        try:
            font_bold = first_run.font.bold
        except:
            pass
        try:
            font_italic = first_run.font.italic
        except:
            pass
        try:
            if first_run.font.color and hasattr(first_run.font.color, 'rgb') and first_run.font.color.rgb:
                font_color_rgb = first_run.font.color.rgb
        except:
            pass
        try:
            para_alignment = first_para.alignment
        except:
            pass
        try:
            para_level = first_para.level if first_para.level is not None else 0
        except:
            para_level = 0
    
    # Clear existing text
    text_frame.clear()
    
    # Split text by lines
    lines = formatted_text.split('\n')
    
    # Helper function to apply formatting to a run
    def apply_formatting(run, is_bold=False, is_header=False, header_text=""):
        try:
            # Apply Calibri (Body) font for all content, except preserve original for headers if needed
            if is_header:
                # Headers can keep original font or use Calibri
                run.font.name = font_name if font_name else "Calibri"
            else:
                # All body content: Calibri (Body), 11.5pt, regular
                run.font.name = "Calibri"
        except:
            pass
        try:
            if is_header:
                # Headers: Set to 14pt for "Areas of Expertise" and "Education" in left box
                # Check if this is a left box header (Areas of Expertise or Education)
                if header_text.strip() in ["Areas of Expertise", "Education"]:
                    run.font.size = Pt(14)
                elif font_size and font_size is not None:
                    # Other headers: preserve original size
                    run.font.size = font_size
            else:
                # Body content: 11.5pt
                run.font.size = Pt(11.5)
        except:
            pass
        try:
            # Set bold: True for headers only, regular (False) for body content
            if is_bold or is_header:
                run.font.bold = True
            else:
                # Body content: regular (not bold)
                run.font.bold = False
        except:
            pass
        try:
            if font_italic is not None:
                run.font.italic = font_italic
        except:
            pass
        try:
            if font_color_rgb:
                run.font.color.rgb = font_color_rgb
        except:
            pass
    
    # Process each line
    for i, line in enumerate(lines):
        if i == 0:
            # First paragraph (reuse existing or create new)
            p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
        else:
            # New paragraph for each line
            p = text_frame.add_paragraph()
        
        # Check if this line is a header (Areas of Expertise, Education, Profile Summary, or Projects)
        is_header = line.strip() in ["Areas of Expertise", "Education", "Profile Summary", "Projects"]
        
        # Set text
        p.text = line
        
        # Apply formatting to runs
        if p.runs:
            apply_formatting(p.runs[0], is_bold=is_header, is_header=is_header, header_text=line)
        
        # Preserve paragraph alignment and level
        try:
            if para_alignment is not None:
                p.alignment = para_alignment
        except:
            pass
        try:
            if para_level is not None:
                p.level = para_level
        except:
            pass




def _format_right_box_content(profile_summary, project1, project2, project3=None, project4=None):
    """Format right box content with Profile Summary and Projects sections
    
    Args:
        profile_summary: Profile summary text string
        project1: Dictionary with project1 information (title, duration, description, technologies)
        project2: Dictionary with project2 information (title, duration, description, technologies)
        project3: Dictionary with project3 information (title, duration, description, technologies) - optional
        project4: Dictionary with project4 information (title, duration, description, technologies) - optional
    
    Returns:
        Formatted text string with headers and content
    """
    formatted_parts = []
    
    # Format Profile Summary section
    if profile_summary and profile_summary.strip() and profile_summary != 'Not Found':
        formatted_parts.append("Profile Summary")
        # Add summary text (can be multiple sentences/paragraphs)
        formatted_parts.append(profile_summary.strip())
    
    # Collect all available projects
    all_projects = []
    if project1:
        all_projects.append(project1)
    if project2:
        all_projects.append(project2)
    if project3:
        all_projects.append(project3)
    if project4:
        all_projects.append(project4)
    
    # Filter out projects with no meaningful data
    valid_projects = []
    for proj in all_projects:
        has_title = proj.get('title', 'Not Found') not in ['Not Found', '', None]
        has_description = proj.get('description', 'Not Found') not in ['Not Found', '', None]
        has_technologies = proj.get('technologies', 'Not Found') not in ['Not Found', '', None]
        if has_title or has_description or has_technologies:
            valid_projects.append(proj)
    
    # Estimate content size and limit projects if needed
    # If we have 4 projects, estimate total content size
    max_projects = 4
    if len(valid_projects) == 4:
        # Estimate content size (character count)
        estimated_size = len(profile_summary) if profile_summary else 0
        for proj in valid_projects:
            estimated_size += len(proj.get('title', '')) + len(proj.get('description', '')) + len(proj.get('technologies', ''))
        
        # If estimated size exceeds threshold (2500 chars), limit to 3 projects
        if estimated_size > 2500:
            max_projects = 3
            valid_projects = valid_projects[:3]
    
    # Format Projects section
    projects_added = False
    project_number = 1
    
    # Add projects up to max_projects limit
    for proj in valid_projects[:max_projects]:
        if not projects_added:
            # Add spacing before Projects section if Profile Summary exists
            if formatted_parts:
                formatted_parts.append("")  # Empty line separator
            formatted_parts.append("Projects")
            projects_added = True
        
        # Add spacing between projects (except for first project)
        if projects_added and formatted_parts and formatted_parts[-1] != "Projects":
            formatted_parts.append("")  # Empty line between projects
        
        # Format project with label
        project_text = []
        project_text.append(f"Project {project_number}")
        
        has_description = proj.get('description', 'Not Found') not in ['Not Found', '', None]
        has_technologies = proj.get('technologies', 'Not Found') not in ['Not Found', '', None]
        
        if has_description:
            # Description is already in bullet format from LLM (no conversion needed)
            desc = proj.get('description', '').strip()
            # Ensure each bullet is on a new line (split by \n if not already)
            # LLM should already format with \n, but ensure it's properly split
            if '\n' in desc:
                # Already has newlines - add each line separately
                desc_lines = desc.split('\n')
                for desc_line in desc_lines:
                    if desc_line.strip():
                        project_text.append(desc_line.strip())
            else:
                # Single line - add as-is
                project_text.append(desc)
        
        if has_technologies:
            project_text.append(f"• Technologies: {proj.get('technologies', '')}")
        
        formatted_parts.extend(project_text)
        project_number += 1
    
    # Join all parts with newlines
    return '\n'.join(formatted_parts)


def _replace_text_preserve_formatting(text_frame, new_text):
    """Replace text in text_frame while preserving original font, size, and style
    
    Args:
        text_frame: TextFrame object from python-pptx
        new_text: New text to insert
    """
    if not text_frame.paragraphs:
        text_frame.add_paragraph()
    
    # Get formatting from first paragraph (preserve font, size, style)
    first_para = text_frame.paragraphs[0]
    
    # Store original formatting
    original_runs = first_para.runs
    font_name = None
    font_size = None
    font_bold = None
    font_italic = None
    font_color_rgb = None
    para_alignment = None
    para_level = 0
    
    if original_runs:
        first_run = original_runs[0]
        try:
            font_name = first_run.font.name if first_run.font.name else None
        except:
            pass
        try:
            font_size = first_run.font.size
        except:
            pass
        try:
            font_bold = first_run.font.bold
        except:
            pass
        try:
            font_italic = first_run.font.italic
        except:
            pass
        try:
            # Font color is complex - try to get RGB value
            if first_run.font.color and hasattr(first_run.font.color, 'rgb') and first_run.font.color.rgb:
                font_color_rgb = first_run.font.color.rgb
        except:
            pass
        try:
            # Paragraph-level formatting
            para_alignment = first_para.alignment
        except:
            pass
        try:
            para_level = first_para.level if first_para.level is not None else 0
        except:
            para_level = 0
    
    # Clear existing text but keep paragraph structure
    text_frame.clear()
    
    # Split new text by newlines to preserve paragraph structure
    lines = new_text.split('\n')
    
    # Helper function to apply formatting to a run
    def apply_formatting(run):
        try:
            if font_name:
                run.font.name = font_name
        except:
            pass
        try:
            if font_size and font_size is not None:
                run.font.size = font_size
        except:
            pass
        try:
            if font_bold is not None:
                run.font.bold = font_bold
        except:
            pass
        try:
            if font_italic is not None:
                run.font.italic = font_italic
        except:
            pass
        try:
            if font_color_rgb:
                run.font.color.rgb = font_color_rgb
        except:
            pass
    
    # Add first paragraph with preserved formatting
    p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    first_line = lines[0] if lines else new_text
    p.text = first_line
    
    # Apply formatting to first run
    if p.runs:
        apply_formatting(p.runs[0])
    
    # Preserve paragraph alignment (only if valid)
    try:
        if para_alignment is not None:
            p.alignment = para_alignment
    except:
        pass
    try:
        if para_level is not None:
            p.level = para_level
    except:
        pass
    
    # Add remaining lines as new paragraphs with same formatting
    for line in lines[1:]:
        if line.strip():  # Only add non-empty lines
            new_p = text_frame.add_paragraph()
            new_p.text = line
            if new_p.runs:
                apply_formatting(new_p.runs[0])
            # Preserve paragraph alignment
            try:
                if para_alignment is not None:
                    new_p.alignment = para_alignment
            except:
                pass
            try:
                if para_level is not None:
                    new_p.level = para_level
            except:
                pass


def read_sample_ppt_structure(sample_ppt_path):
    """Read sample PPT and return structure information
    
    Args:
        sample_ppt_path: Path to sample PPT file
    
    Returns:
        Presentation object
    """
    try:
        prs = Presentation(sample_ppt_path)
        return prs
    except Exception as e:
        raise Exception(f"Error reading sample PPT: {str(e)}")


def create_ppt_from_sample(sample_ppt_path, candidate_data, output_path):
    """Create new PPT from sample with candidate data
    
    Args:
        sample_ppt_path: Path to sample PPT file
        candidate_data: Dictionary with candidate information:
            - candidate_name: Candidate name
            - position: Position/Job title
            - location: Location
            - area_of_expertise: Technical skills (comma-separated)
            - education: Education details
            - profile_summary: Profile summary
            - project1: Dictionary with {title, duration, description, technologies}
            - project2: Dictionary with {title, duration, description, technologies}
        output_path: Path to save the new PPT
    
    Returns:
        True if successful, False otherwise
    """
    try:
        # Load sample presentation
        prs = Presentation(sample_ppt_path)
        
        # Get the first slide (assuming single slide template)
        if len(prs.slides) == 0:
            return False
        
        slide = prs.slides[0]
        
        # Identify text boxes by position: upper, left, right
        text_boxes = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text_boxes.append({
                    'shape': shape,
                    'text_frame': shape.text_frame,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                })
        
        # Store references to upper and left boxes for later use
        upper_box_ref = None
        left_box_ref = None
        right_box_ref = None
        
        # Sort to identify: upper (lowest top), left (lowest left), right (highest left)
        if text_boxes:
            # Find upper box (lowest top value)
            upper_box = min(text_boxes, key=lambda x: x['top'])
            upper_box_ref = upper_box
            
            # Find left and right boxes (excluding upper box)
            side_boxes = [box for box in text_boxes if box != upper_box]
            if side_boxes:
                left_box = min(side_boxes, key=lambda x: x['left'])
                right_box = max(side_boxes, key=lambda x: x['left'])
                left_box_ref = left_box
                right_box_ref = right_box
            
            # Fill upper box: Candidate_Name | Position | Location
            if upper_box:
                upper_text = []
                if candidate_data.get('candidate_name') and candidate_data.get('candidate_name') != 'Not Found':
                    upper_text.append(candidate_data['candidate_name'])
                if candidate_data.get('position') and candidate_data.get('position') != 'Not Found':
                    upper_text.append(candidate_data['position'])
                if candidate_data.get('location') and candidate_data.get('location') != 'Not Found':
                    upper_text.append(candidate_data['location'])
                
                if upper_text:
                    combined_upper = ' | '.join(upper_text)
                    _replace_text_preserve_formatting(upper_box['text_frame'], combined_upper)
            
            # Fill left box: Areas of Expertise + Education (formatted with headers and bullets)
            if left_box:
                formatted_left_text = _format_left_box_content(
                    candidate_data.get('area_of_expertise', ''),
                    candidate_data.get('education', '')
                )
                
                if formatted_left_text:
                    _replace_text_with_bold_headers(left_box['text_frame'], formatted_left_text)
        
            # Fill right box: Profile Summary + Projects (formatted with headers and bullets)
            if right_box:
                formatted_right_text = _format_right_box_content(
                    candidate_data.get('profile_summary', ''),
                    candidate_data.get('project1', {}),
                    candidate_data.get('project2', {}),
                    candidate_data.get('project3', {}),
                    candidate_data.get('project4', {})
                )
                
                if formatted_right_text:
                    _replace_text_with_bold_headers(right_box['text_frame'], formatted_right_text)
        
        # Save the presentation
        prs.save(output_path)
        return True
        
    except Exception as e:
        raise Exception(f"Error creating PPT: {str(e)}")



